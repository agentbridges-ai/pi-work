#!/usr/bin/env python3
"""Generate the deterministic PowerPoint-family preview showcase source."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.opc.package import Part
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "demo-user-space" / "office"
OUTPUT = OUT_DIR / "ONLYOFFICE Presentation Preview Showcase.pptx"
ICON = ROOT / "web" / "public" / "icons" / "piwork-512.png"

NAVY = RGBColor(0x12, 0x30, 0x4A)
BLUE = RGBColor(0x2F, 0x75, 0xB5)
TEAL = RGBColor(0x00, 0xA6, 0xA6)
GOLD = RGBColor(0xF2, 0xB1, 0x34)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
PALE = RGBColor(0xEA, 0xF3, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x60, 0x77, 0x85)


def set_background(slide, color=NAVY) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, size=20, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(x, y, w, h)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Pt(4)
    frame.margin_right = Pt(4)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_title(slide, title: str, marker: str) -> None:
    add_text(slide, Inches(0.65), Inches(0.25), Inches(11.9), Inches(0.55), title, 25, NAVY, True)
    add_text(slide, Inches(0.68), Inches(0.82), Inches(11.7), Inches(0.3), marker, 9, CORAL, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.13), Inches(11.75), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_footer(slide, number: int) -> None:
    add_text(slide, Inches(0.65), Inches(7.18), Inches(7.0), Inches(0.22), "ONLYOFFICE · deterministic preview corpus", 8, MUTED)
    add_text(slide, Inches(11.9), Inches(7.14), Inches(0.5), Inches(0.25), str(number), 9, MUTED, True, PP_ALIGN.RIGHT)


def add_transition(slide, kind: str = "fade") -> None:
    slide_el = slide._element
    transition = parse_xml(f'<p:transition {nsdecls("p")} spd="med"><p:{kind}/></p:transition>')
    timing = slide_el.find(qn("p:timing"))
    if timing is not None:
        timing.addprevious(transition)
    else:
        ext_lst = slide_el.find(qn("p:extLst"))
        if ext_lst is not None:
            ext_lst.addprevious(transition)
        else:
            slide_el.append(transition)


def add_basic_timing(slide, shape_id: int) -> None:
    timing = parse_xml(
        f'<p:timing {nsdecls("p")}><p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        '<p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        '<p:childTnLst><p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        '<p:childTnLst><p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="4" dur="700" fill="hold"/>'
        f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl></p:cBhvr></p:animEffect></p:childTnLst>'
        '</p:cTn></p:par></p:childTnLst></p:cTn></p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
    )
    slide._element.append(timing)


def add_gradient(shape, start: str, end: str) -> None:
    sp_pr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        node = sp_pr.find(qn(tag))
        if node is not None:
            sp_pr.remove(node)
    gradient = parse_xml(
        f'<a:gradFill {nsdecls("a")} rotWithShape="1"><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{start}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{end}"/></a:gs>'
        '</a:gsLst><a:lin ang="2700000" scaled="1"/></a:gradFill>'
    )
    xfrm = sp_pr.find(qn("a:xfrm"))
    sp_pr.insert(sp_pr.index(xfrm) + 1 if xfrm is not None else 0, gradient)


def add_shadow(shape) -> None:
    sp_pr = shape._element.spPr
    effect = parse_xml(
        f'<a:effectLst {nsdecls("a")}><a:outerShdw blurRad="50800" dist="38100" dir="2700000" '
        'algn="ctr" rotWithShape="0"><a:srgbClr val="12304A"><a:alpha val="35000"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
    )
    sp_pr.append(effect)


def add_svg(slide, x, y, w, h) -> None:
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="640" height="420" viewBox="0 0 640 420">'
        '<rect width="640" height="420" rx="40" fill="#12304A"/>'
        '<path d="M100 210h420" stroke="#F2B134" stroke-width="28" stroke-linecap="round"/>'
        '<circle cx="160" cy="210" r="78" fill="#00A6A6"/>'
        '<circle cx="480" cy="210" r="78" fill="#2F75B5"/>'
        '<path d="M445 210l-48-42v84z" fill="white"/></svg>'
    ).encode()
    part_name = PackURI(f"/ppt/media/showcase-vector-{len(slide.part.rels)}.svg")
    part = Part(part_name, "image/svg+xml", slide.part.package, svg)
    rel_id = slide.part.relate_to(part, RT.IMAGE)
    pic = parse_xml(
        f'<p:pic {nsdecls("p", "a", "r")}><p:nvPicPr><p:cNvPr id="606" name="P06 SVG vector" '
        'descr="P06 SVG vector with non-zero dimensions"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>'
        f'<p:blipFill><a:blip r:embed="{rel_id}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        f'<p:spPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
    )
    slide.shapes._spTree.insert_element_before(pic, "p:extLst")


def style_chart(chart) -> None:
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
    palette = (BLUE, CORAL, TEAL, GOLD)
    for index, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = palette[index % len(palette)]


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # P01 cover/theme/master geometry
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    band.fill.solid(); band.fill.fore_color.rgb = TEAL; band.line.fill.background()
    add_text(slide, Inches(0.8), Inches(1.05), Inches(10.8), Inches(0.9), "Presentation Preview Showcase", 32, WHITE, True)
    add_text(slide, Inches(0.85), Inches(2.0), Inches(10.6), Inches(0.55), "Theme · masters · text · shapes · media · tables · charts · transitions", 15, PALE)
    if ICON.exists():
        pic = slide.shapes.add_picture(str(ICON), Inches(10.9), Inches(0.85), width=Inches(1.35))
        pic.name = "P06 raster image with alt text"
        pic._element.nvPicPr.cNvPr.set("descr", "P06 raster image with non-zero dimensions")
    for i, (label, color) in enumerate((("TEXT", BLUE), ("DRAW", TEAL), ("DATA", GOLD), ("MOTION", CORAL))):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85 + i * 2.85), Inches(3.15), Inches(2.5), Inches(1.25))
        card.fill.solid(); card.fill.fore_color.rgb = color; card.line.fill.background()
        tf = card.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label; r.font.name = "Aptos"; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE
    add_text(slide, Inches(0.85), Inches(6.55), Inches(10), Inches(0.3), "P01 · 16:9 slide geometry · deterministic theme palette · layout references", 10, PALE)
    add_footer(slide, 1)
    add_transition(slide)

    # P02 rich text and lists
    slide = prs.slides.add_slide(blank); add_title(slide, "Text hierarchy and list semantics", "P02/P03 · distinct text runs, paragraph levels, bullets and numbering")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(5.8), Inches(4.9))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Rich run formatting"; p.level = 0
    p.runs[0].font.size = Pt(24); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    p = tf.add_paragraph();
    for text, attrs in (("Bold", {"bold": True}), (" · italic", {"italic": True}), (" · underline", {"underline": True}), (" · coral", {"color": CORAL})):
        r = p.add_run(); r.text = text; r.font.size = Pt(18); r.font.name = "Aptos"
        if attrs.get("bold"): r.font.bold = True
        if attrs.get("italic"): r.font.italic = True
        if attrs.get("underline"): r.font.underline = True
        r.font.color.rgb = attrs.get("color", NAVY)
    for level, text in ((0, "Primary bullet"), (1, "Nested supporting point"), (1, "Nested verification marker"), (0, "Second primary bullet")):
        p = tf.add_paragraph(); p.text = text; p.level = level; p.font.size = Pt(17); p.font.color.rgb = NAVY
    rtl = add_text(slide, Inches(7.0), Inches(1.65), Inches(5.2), Inches(1.2), "简体中文 · 日本語\nالعربية · עברית", 20, NAVY, False, PP_ALIGN.RIGHT)
    rtl.name = "P02 multilingual text"
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(3.2), Inches(5.2), Inches(1.5))
    quote.fill.solid(); quote.fill.fore_color.rgb = PALE; quote.line.color.rgb = BLUE
    quote.text_frame.text = "Text remains editable, searchable and selectable after browser conversion."
    quote.text_frame.paragraphs[0].font.size = Pt(18); quote.text_frame.paragraphs[0].font.color.rgb = NAVY
    add_footer(slide, 2); add_transition(slide, "push")

    # P04/P05 shapes, gradients, connectors and groups
    slide = prs.slides.add_slide(blank); add_title(slide, "Shapes, fills, connectors and grouping", "P04/P05 · gradient, outline, shadow, connector endpoints, z-order and group model")
    group = slide.shapes.add_group_shape()
    group.name = "P05 grouped process"
    nodes = []
    for i, (label, color) in enumerate((("SOURCE", BLUE), ("CONVERT", TEAL), ("MODEL", GOLD), ("VERIFY", CORAL))):
        node = group.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(i * 2.25), Inches(0), Inches(1.85), Inches(1.0))
        node.fill.solid(); node.fill.fore_color.rgb = color; node.line.color.rgb = WHITE
        node.text_frame.text = label; node.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        node.text_frame.paragraphs[0].font.bold = True; node.text_frame.paragraphs[0].font.color.rgb = WHITE
        nodes.append(node)
    group.left, group.top, group.width, group.height = Inches(0.9), Inches(1.6), Inches(9.0), Inches(1.4)
    for i in range(3):
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.75 + i * 2.25), Inches(2.3), Inches(3.05 + i * 2.25), Inches(2.3))
        connector.line.color.rgb = NAVY; connector.line.width = Pt(2.25)
    shapes = ((MSO_SHAPE.HEXAGON, "GRADIENT", "2F75B5", "00A6A6"), (MSO_SHAPE.CHEVRON, "SHADOW", "F2B134", "E76F51"), (MSO_SHAPE.CLOUD, "OUTLINE", "EAF3F8", "EAF3F8"))
    for i, (kind, label, start, end) in enumerate(shapes):
        shape = slide.shapes.add_shape(kind, Inches(1.0 + i * 4.0), Inches(3.55), Inches(3.1), Inches(1.65))
        shape.text_frame.text = label; shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].font.bold = True; shape.text_frame.paragraphs[0].font.size = Pt(18); shape.text_frame.paragraphs[0].font.color.rgb = NAVY
        add_gradient(shape, start, end); shape.line.color.rgb = NAVY
        if i == 1: add_shadow(shape)
    add_footer(slide, 3); add_transition(slide, "wipe")

    # P06 media
    slide = prs.slides.add_slide(blank); add_title(slide, "Raster and vector media", "P06 · decoded picture resources, crop geometry and alternative text")
    if ICON.exists():
        pic = slide.shapes.add_picture(str(ICON), Inches(0.9), Inches(1.5), width=Inches(3.0))
        pic.name = "P06 primary raster"; pic._element.nvPicPr.cNvPr.set("descr", "P06 primary raster alternative text")
        crop = slide.shapes.add_picture(str(ICON), Inches(4.7), Inches(1.5), width=Inches(3.0), height=Inches(3.0))
        crop.crop_left = 0.12; crop.crop_right = 0.12; crop.crop_top = 0.12; crop.crop_bottom = 0.12
        crop.name = "P06 cropped raster"; crop._element.nvPicPr.cNvPr.set("descr", "P06 cropped raster alternative text")
    add_svg(slide, Inches(8.55), Inches(1.55), Inches(3.3), Inches(2.2))
    add_text(slide, Inches(8.1), Inches(4.15), Inches(4.0), Inches(0.8), "P06 SVG resource\nremains resolution-independent", 17, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 4); add_transition(slide)

    # P07 table
    slide = prs.slides.add_slide(blank); add_title(slide, "Table grid, merge and cell styling", "P07 · 4×4 table, merged region, header fill, borders and alignment")
    table = slide.shapes.add_table(4, 4, Inches(0.85), Inches(1.45), Inches(11.65), Inches(3.45)).table
    labels = (("Marker", "Element", "Expected model", "Status"), ("P07-A", "Header row", "fill + bold", "READY"), ("P07-B", "Merged region", "gridSpan + borders", "READY"), ("P07-C", "Cell alignment", "center/middle", "READY"))
    for r, values in enumerate(labels):
        for c, value in enumerate(values):
            cell = table.cell(r, c); cell.text = value; cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.text_frame.paragraphs[0].font.size = Pt(15); cell.text_frame.paragraphs[0].font.color.rgb = WHITE if r == 0 else NAVY
            cell.text_frame.paragraphs[0].font.bold = r == 0 or c == 0
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY if r == 0 else (PALE if r % 2 else WHITE)
    table.cell(2, 1).merge(table.cell(2, 2)); table.cell(2, 1).text = "P07-B merged across two cells"
    add_footer(slide, 5); add_transition(slide, "cover")

    # P08 charts
    slide = prs.slides.add_slide(blank); add_title(slide, "Native chart objects", "P08 · column, line and pie chart types with multiple colored series")
    data = CategoryChartData(); data.categories = ["Jan", "Feb", "Mar", "Apr"]
    data.add_series("Document", (42, 57, 66, 81)); data.add_series("Spreadsheet", (35, 49, 61, 72))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.4), Inches(4.1), Inches(4.9), data).chart
    chart.has_title = True; chart.chart_title.text_frame.text = "Column"; style_chart(chart)
    line_data = CategoryChartData(); line_data.categories = ["Jan", "Feb", "Mar", "Apr"]
    line_data.add_series("Preview", (68, 74, 83, 92)); line_data.add_series("Target", (72, 78, 85, 94))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(4.65), Inches(1.4), Inches(4.1), Inches(4.9), line_data).chart
    chart.has_title = True; chart.chart_title.text_frame.text = "Line"; style_chart(chart)
    pie_data = CategoryChartData(); pie_data.categories = ["Text", "Draw", "Data", "Motion"]; pie_data.add_series("Coverage", (38, 27, 24, 11))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(8.75), Inches(1.4), Inches(4.0), Inches(4.9), pie_data).chart
    chart.has_title = True; chart.chart_title.text_frame.text = "Pie"; style_chart(chart)
    add_footer(slide, 6); add_transition(slide)

    # P09 links/equation/symbols
    slide = prs.slides.add_slide(blank); add_title(slide, "Links, symbols and equation content", "P09/P10 · equation marker, URL, mailto and internal slide action")
    eq = add_text(slide, Inches(1.0), Inches(1.6), Inches(11.2), Inches(1.2), "E = mc²     ∫₀¹ x² dx = ⅓     Ω · ∑ · √", 30, NAVY, True, PP_ALIGN.CENTER)
    url = add_text(slide, Inches(1.2), Inches(3.2), Inches(5.0), Inches(0.6), "ONLYOFFICE DocumentServer", 19, BLUE, True)
    url.text_frame.paragraphs[0].runs[0].hyperlink.address = "https://github.com/ONLYOFFICE/DocumentServer"
    mail = add_text(slide, Inches(1.2), Inches(4.1), Inches(5.0), Inches(0.6), "preview@example.test", 19, TEAL, True)
    mail.text_frame.paragraphs[0].runs[0].hyperlink.address = "mailto:preview@example.test"
    internal = slide.shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_BACK_OR_PREVIOUS, Inches(8.8), Inches(3.0), Inches(1.2), Inches(1.2))
    internal.fill.solid(); internal.fill.fore_color.rgb = GOLD; internal.click_action.target_slide = prs.slides[0]
    add_text(slide, Inches(7.8), Inches(4.4), Inches(3.2), Inches(0.5), "Internal link → slide 1", 15, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 7); add_transition(slide, "split")

    # P11/P12 notes, transition and animation
    slide = prs.slides.add_slide(blank); add_title(slide, "Notes, placeholders and motion metadata", "P11/P12 · speaker notes, footer/date/number markers, transition and timing tree")
    motion = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.8), Inches(4.2), Inches(2.3))
    motion.fill.solid(); motion.fill.fore_color.rgb = BLUE; motion.line.fill.background(); motion.text_frame.text = "P12 ANIMATED OBJECT"
    motion.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER; motion.text_frame.paragraphs[0].font.size = Pt(22); motion.text_frame.paragraphs[0].font.bold = True; motion.text_frame.paragraphs[0].font.color.rgb = WHITE
    add_basic_timing(slide, motion.shape_id)
    notes = slide.notes_slide.notes_text_frame
    notes.text = "P11 speaker notes: this text must remain in notesSlide resources and the editor notes model."
    add_text(slide, Inches(6.4), Inches(1.8), Inches(5.2), Inches(2.4), "Transition: fade\nAnimation: entrance/fade\nNotes: notesSlide XML\nFooter/date/number: stable markers", 19, NAVY, False)
    add_footer(slide, 8); add_transition(slide)

    # P13 SmartArt-compatible grouped diagram fallback
    slide = prs.slides.add_slide(blank); add_title(slide, "Diagram and group-object preview", "P13 · grouped process diagram retained as editable vector shapes")
    group = slide.shapes.add_group_shape(); group.name = "P13 grouped diagram"
    for i, (label, color) in enumerate((("LOAD", BLUE), ("X2T", TEAL), ("VIEW", GOLD), ("TEST", CORAL))):
        node = group.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(i * 2.25), Inches(0), Inches(2.5), Inches(1.55))
        node.fill.solid(); node.fill.fore_color.rgb = color; node.line.color.rgb = WHITE
        node.text_frame.text = label; node.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        node.text_frame.paragraphs[0].font.size = Pt(17); node.text_frame.paragraphs[0].font.bold = True; node.text_frame.paragraphs[0].font.color.rgb = WHITE
    group.left, group.top, group.width, group.height = Inches(1.0), Inches(2.1), Inches(11.3), Inches(2.2)
    add_text(slide, Inches(1.2), Inches(5.1), Inches(10.9), Inches(0.8), "P13 browser acceptance inspects the group and child drawing objects, not screenshot pixels.", 18, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 9); add_transition(slide, "push")

    prs.core_properties.title = "ONLYOFFICE Presentation Preview Showcase"
    prs.core_properties.subject = "Deterministic presentation rendering and document-model fixture"
    prs.core_properties.author = "Piwork"
    prs.core_properties.keywords = "ONLYOFFICE, PPTX, preview, chart, transition, animation"
    prs.core_properties.comments = "Generated by scripts/generate-office-presentation-showcase.py"
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
